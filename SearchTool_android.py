# --- START OF FILE SearchTool.py ---

import os
import re
import fitz  # PyMuPDF for PDFs
import docx
import pptx
import pandas as pd
import subprocess
# import platform # Replaced by kivy.utils.platform check below
from kivy.app import App
from kivy.uix.boxlayout import BoxLayout
from kivy.uix.gridlayout import GridLayout
from kivy.uix.button import Button
from kivy.uix.label import Label
from kivy.uix.textinput import TextInput
from kivy.uix.scrollview import ScrollView
# from kivy.uix.filechooser import FileChooserListView # Removed
from kivy.uix.popup import Popup
from kivy.core.window import Window
from kivy.uix.behaviors import ButtonBehavior
from kivy.uix.modalview import ModalView
from kivy.graphics import Color, Rectangle
from kivy.utils import get_color_from_hex, platform # Import platform
from kivy.clock import Clock
from functools import partial

# ---vvv--- ADDED PLYER IMPORTS ---vvv---
try:
    from plyer import filechooser, permissions, fileopener # Import necessary plyer modules
except ImportError:
    print("WARNING: Plyer library not found. Install it with 'pip install plyer'. File/Folder choosing and opening will likely fail on Android.")
    # Define dummy functions or classes if plyer is missing, so the app doesn't crash immediately
    # This is optional but can help basic UI testing without full functionality
    class DummyFileChooser:
        def choose_dir(self, **kwargs): print("Plyer not found: choose_dir unavailable")
        def open_file(self, **kwargs): print("Plyer not found: open_file unavailable")
    class DummyPermissions:
        Permission = type('obj', (object,), {'READ_EXTERNAL_STORAGE': 'READ_EXTERNAL_STORAGE', 'WRITE_EXTERNAL_STORAGE': 'WRITE_EXTERNAL_STORAGE'})()
        def request(self, perms, callback): print("Plyer not found: permissions.request unavailable"); callback(perms, [False]*len(perms)) if callback else None
    class DummyFileOpener:
        def open(self, path): print(f"Plyer not found: Cannot open {path}")

    filechooser = DummyFileChooser()
    permissions = DummyPermissions()
    fileopener = DummyFileOpener()
# ---^^^--- END PLYER IMPORTS ---^^^---


# Set theme colors
THEME = {
    'primary': '#738c63',  # Blue
    'secondary': '#A0B38C',  # Light blue
    'text': '#333333',  # Dark gray
    'background': '#A0B38C',  # Off-white
    'accent': '#ff6b6b',  # Accent red
    'success': '#5cb85c',  # Green
    'hover': '#d9e6f7'  # Hover light blue
}
RESULT_ITEM_COLORS = {
    'normal': '#b0c49a',  # Lighter Green (Slightly lighter than background)
    'hover': '#A0B38C',   # Main Background Green (Theme secondary)
    'press': '#889877'    # Darker Green
}

# Function definitions (extract_text_from_..., boolean_search) remain the same
# ... (Keep all extract_text functions and boolean_search as they were) ...
def extract_text_from_pdf(filepath):
    text = ""
    try:
        doc = fitz.open(filepath)
        for page in doc:
            text += page.get_text("text") + "\n"
    except Exception as e:
        print(f"Error reading {filepath}: {e}")
    return text

def extract_text_from_docx(filepath):
    try:
        doc = docx.Document(filepath)
        return "\n".join([para.text for para in doc.paragraphs])
    except Exception as e:
        print(f"Error reading {filepath}: {e}")
        return ""

def extract_text_from_pptx(filepath):
    try:
        presentation = pptx.Presentation(filepath)
        text = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)
    except Exception as e:
        print(f"Error reading {filepath}: {e}")
        return ""

def extract_text_from_excel(filepath):
    try:
        df = pd.read_excel(filepath, sheet_name=None)
        text = "\n".join([df[sheet].to_string() for sheet in df])
        return text
    except Exception as e:
        print(f"Error reading {filepath}: {e}")
        return ""

def extract_text_from_txt(filepath):
    try:
        with open(filepath, 'r', encoding='utf-8', errors='ignore') as file:
            return file.read()
    except Exception as e:
        print(f"Error reading {filepath}: {e}")
        return ""

def extract_text_from_csv(filepath):
    try:
        df = pd.read_csv(filepath)
        return df.to_string()
    except Exception as e:
        print(f"Error reading {filepath}: {e}")
        return ""

def extract_text(filepath):
    ext = filepath.split(".")[-1].lower()
    if ext == "pdf":
        return extract_text_from_pdf(filepath)
    elif ext == "docx":
        return extract_text_from_docx(filepath)
    elif ext == "pptx":
        return extract_text_from_pptx(filepath)
    elif ext in ["xls", "xlsx"]:
        return extract_text_from_excel(filepath)
    elif ext == "txt":
        return extract_text_from_txt(filepath)
    elif ext == "csv":
        return extract_text_from_csv(filepath)
    elif ext == "rtf":
        try:
            # Basic RTF handling (same as before)
            with open(filepath, 'r', encoding='utf-8', errors='ignore') as file:
                content = file.read()
                content = re.sub(r'\{\*?\\[^{}]+}|[{}]|\\\w+(\s?)|\\.\s?', '', content)
                return content.strip()
        except Exception as e:
            print(f"Error reading RTF {filepath}: {e}")
            return extract_text_from_txt(filepath) # Fallback
    elif ext in ["json", "xml", "html", "htm", "md", "log"]:
        return extract_text_from_txt(filepath)
    else:
        print(f"Unsupported format: {filepath}")
        return "" # Return empty string for unsupported or failed extractions

def boolean_search(text, query, exact_match=False):
    # Using the more robust version from previous examples
    processed_query = query.replace(" AND ", " & ").replace(" OR ", " | ").replace(" NOT ", " ~")
    terms = re.findall(r'"[^"]+"|\b\w+\b', query)
    term_results = {}

    for term in terms:
        original_term = term
        is_phrase = term.startswith('"') and term.endswith('"')
        if is_phrase:
            term = term[1:-1]

        if exact_match:
            pattern = r'\b' + re.escape(term) + r'\b'
        else:
            pattern = re.escape(term)

        found = bool(re.search(pattern, text, re.IGNORECASE))
        term_results[original_term] = found

    sorted_terms = sorted(term_results.keys(), key=len, reverse=True)
    eval_string = processed_query
    for term in sorted_terms:
        # Replace term with its boolean value (True/False)
        # Need careful replacement. Using regex to ensure whole word/phrase replacement might be safer,
        # but simple replace often works if query structure is reasonable.
        # Let's refine this slightly with word boundaries for standalone words
        if term.startswith('"'):
            # Replace quoted phrase directly
            eval_string = eval_string.replace(term, str(term_results[term]))
        else:
            # Replace standalone word using regex for boundary matching
            # This handles cases like 'on' vs 'ontology' better.
            # Pattern explanation:
            # (?<![\w&|~])  # Negative lookbehind: Not preceded by word char or operator symbol
            # \b{term}\b   # The word itself with boundaries
            # (?![\w&|~])   # Negative lookahead: Not followed by word char or operator symbol
            # Note: This lookahead/behind might be too strict if operators are immediately next to words.
            # Let's revert to simpler replacement for now, assuming spaces around terms/operators.
            # If issues arise, refine the replacement logic.
             eval_string = eval_string.replace(term, str(term_results[term]))


    eval_string = eval_string.replace('&', ' and ').replace('|', ' or ').replace('~', ' not ')

    try:
        # Basic validation
        if not re.fullmatch(r'^[() TrueFalsenotandor\s]+$', eval_string):
             print(f"Warning: Invalid characters in eval string: {eval_string}")
             return False
        # Evaluate
        return bool(eval(eval_string, {'__builtins__': {}}, {'True': True, 'False': False}))
    except Exception as e:
        print(f"Error evaluating boolean query '{eval_string}' derived from '{query}': {e}")
        return False


# Themed Widgets (ThemedButton, ThemedLabel, ThemedTextInput) remain the same
# ... (Keep these classes as they were) ...
class ThemedButton(Button):
    def __init__(self, **kwargs):
        super(ThemedButton, self).__init__(**kwargs)
        self.background_normal = ''
        self.background_color = get_color_from_hex(THEME['primary'])
        self.color = get_color_from_hex('#ffffff')  # White text
        self.border = (0, 0, 0, 0)

    def on_press(self):
        # Use a slightly darker shade for pressed state if needed
        pass # Keep primary color on press for simplicity, or adjust

    def on_release(self):
        # Restore original color on release
        self.background_color = get_color_from_hex(THEME['primary'])

class ThemedLabel(Label):
    def __init__(self, **kwargs):
        super(ThemedLabel, self).__init__(**kwargs)
        self.color = get_color_from_hex(THEME['text'])

class ThemedTextInput(TextInput):
    def __init__(self, **kwargs):
        super(ThemedTextInput, self).__init__(**kwargs)
        self.background_color = get_color_from_hex('#ffffff')
        self.foreground_color = get_color_from_hex(THEME['text'])
        self.cursor_color = get_color_from_hex(THEME['primary'])
        self.selection_color = get_color_from_hex(THEME['secondary'])
        self.padding = [10, 10, 10, 10]
        self.cursor_width = 2


# ResultItem class (with hover effects and MODIFIED open_file)
class ResultItem(ButtonBehavior, BoxLayout):
    def __init__(self, filename, **kwargs):
        super(ResultItem, self).__init__(**kwargs)
        self.orientation = 'horizontal'
        self.size_hint_y = None
        self.height = 50
        self.filename = filename
        self.last_click = 0
        self.is_hovering = False

        self.normal_color_hex = RESULT_ITEM_COLORS['normal']
        self.hover_color_hex = RESULT_ITEM_COLORS['hover']
        self.press_color_hex = RESULT_ITEM_COLORS['press']

        with self.canvas.before:
            self.bg_color = Color(*get_color_from_hex(self.normal_color_hex))
            self.bg_rect = Rectangle(pos=self.pos, size=self.size)

        self.bind(pos=self.update_rect, size=self.update_rect)

        self.label = ThemedLabel(
            text=filename,
            size_hint_x=0.9,
            halign='left',
            valign='middle',
            shorten=True,
            shorten_from='right'
        )
        self.label.bind(size=self.label.setter('text_size'))
        self.add_widget(self.label)

        self.bind(on_press=self.on_item_press)
        self.bind(on_release=self.on_item_release)

        # Only bind mouse hover events if not on Android/iOS (touch platforms)
        if platform not in ('android', 'ios'):
            Window.bind(mouse_pos=self.on_mouse_pos)

    def update_rect(self, *args):
        self.bg_rect.pos = self.pos
        self.bg_rect.size = self.size
        if hasattr(self, 'label'):
             self.label.text_size = (self.width * 0.9, None)

    def on_mouse_pos(self, window, pos):
         # This method should only be called if not on android/ios due to the bind condition
        widget_pos = self.to_widget(*pos)
        is_inside = self.collide_point(*widget_pos)

        if self.is_hovering != is_inside:
            self.is_hovering = is_inside
            if is_inside:
                if self.state == 'normal':
                    self.on_enter()
            else:
                if self.state == 'normal':
                    self.on_leave()

    def on_enter(self):
        # This method is primarily for desktop hover effect
        Window.set_system_cursor('hand')
        self.bg_color.rgba = get_color_from_hex(self.hover_color_hex)

    def on_leave(self):
        # This method is primarily for desktop hover effect
        Window.set_system_cursor('arrow')
        self.bg_color.rgba = get_color_from_hex(self.normal_color_hex)

    def on_item_press(self, instance):
        self.bg_color.rgba = get_color_from_hex(self.press_color_hex)
        current_time = Clock.get_time()
        if current_time - self.last_click < 0.3:
            self.open_file()
        self.last_click = current_time

    def on_item_release(self, instance):
        # Simplified release for touch: always go back to normal color
        # For desktop, check hover state
        if platform not in ('android', 'ios'):
            current_mouse_pos = Window.mouse_pos
            widget_pos = self.to_widget(*current_mouse_pos)
            is_still_inside = self.collide_point(*widget_pos)

            if is_still_inside:
                self.is_hovering = True
                self.bg_color.rgba = get_color_from_hex(self.hover_color_hex)
                Window.set_system_cursor('hand')
            else:
                self.is_hovering = False
                self.bg_color.rgba = get_color_from_hex(self.normal_color_hex)
                Window.set_system_cursor('arrow')
        else:
            # On touch devices, just revert to normal color on release
            self.bg_color.rgba = get_color_from_hex(self.normal_color_hex)


    # ---vvv--- MODIFIED open_file METHOD using Plyer ---vvv---
    def open_file(self):
        app = App.get_running_app()
        if not app or not hasattr(app, 'resume_folder') or not app.resume_folder:
            print("Error: App instance or resume_folder not found/set.")
            ErrorPopup(message="Cannot open file: Folder not set.").open()
            return

        try:
            filepath = os.path.join(app.resume_folder, self.filename)
            print(f"Attempting to open file: {filepath}") # Debug print

            if not os.path.exists(filepath):
                print(f"Error: File not found at {filepath}")
                ErrorPopup(message=f"File not found:\n{self.filename}").open()
                return

            # Use plyer.fileopener for cross-platform opening (including Android Intents)
            try:
                fileopener.open(path=filepath) # Use fileopener instead of platform checks
            except Exception as opener_e:
                 # Fallback for desktop if plyer fails or isn't fully implemented
                 print(f"Plyer fileopener failed: {opener_e}. Falling back for desktop.")
                 current_platform = platform # Get kivy's platform identifier
                 if current_platform == "win":
                     os.startfile(filepath)
                 elif current_platform == "macosx": # Kivy uses 'macosx' for macOS
                     subprocess.run(["open", filepath], check=True)
                 elif current_platform == "linux":
                     subprocess.run(["xdg-open", filepath], check=True)
                 else:
                      # Check for other common Kivy platform strings if necessary
                      print(f"Unsupported platform for fallback open: {current_platform}")
                      raise Exception(f"Unsupported platform for fallback open: {current_platform}")

        except Exception as e:
            print(f"Error opening file '{self.filename}': {e}")
            import traceback
            traceback.print_exc() # Print full traceback for debugging
            ErrorPopup(message=f"Could not open file:\n{e}").open()
    # ---^^^--- END MODIFIED open_file ---^^^---


# ---vvv--- DELETED CustomFileChooserListView CLASS ---vvv---
# ---vvv--- DELETED FolderChooserPopup CLASS ---vvv---


# ErrorPopup class remains the same
# ... (Keep this class as it was) ...
class ErrorPopup(Popup):
    def __init__(self, message, **kwargs):
        super(ErrorPopup, self).__init__(**kwargs)
        self.title = "Error"
        self.size_hint = (None, None)
        self.size = (Window.width * 0.8, Window.height * 0.4) # Adjust size
        self.auto_dismiss = False # Require explicit dismissal for errors
        self.background = ''
        self.background_color = get_color_from_hex('#ffffff')
        self.title_color = get_color_from_hex(THEME['accent'])
        self.separator_color = get_color_from_hex(THEME['accent'])

        content = BoxLayout(orientation='vertical', spacing=10, padding=10)
        # Allow message label to wrap text and scroll if needed
        scroll_view = ScrollView(size_hint_y=0.8)
        msg_label = ThemedLabel(text=message, size_hint_y=None, markup=True) # Enable markup if needed
        msg_label.bind(
             texture_size=lambda instance, size: setattr(instance, 'height', size[1]),
             width=lambda instance, width: setattr(instance, 'text_size', (width, None))
        )
        scroll_view.add_widget(msg_label)
        content.add_widget(scroll_view)

        btn = ThemedButton(text="OK", size_hint_y=None, height=50)
        btn.bind(on_release=self.dismiss)
        content.add_widget(btn)

        self.content = content


# Main App Class with Android modifications
class ResumeSearchApp(App):

    # ---vvv--- ADDED on_start METHOD ---vvv---
    def on_start(self):
        """
        Request necessary permissions on Android when the app starts.
        """
        if platform == 'android':
            print("Requesting Android permissions...")
            try:
                permissions.request(
                    [
                        permissions.Permission.READ_EXTERNAL_STORAGE,
                        permissions.Permission.WRITE_EXTERNAL_STORAGE # Often needed together
                    ],
                    self.permission_callback # Pass a callback function
                )
            except Exception as e:
                 print(f"Error requesting permissions: {e}")
                 ErrorPopup(message=f"Error requesting permissions:\n{e}\nFile access may not work.").open()
        else:
            print("Not Android, skipping permission request.")

    def permission_callback(self, permission_list, grant_result_list):
        """
        Callback function after permissions are requested.
        """
        print("Permission callback received:")
        print(f"  Permissions: {permission_list}")
        print(f"  Grant Results: {grant_result_list}")
        if not all(grant_result_list):
             print("WARNING: Not all storage permissions were granted. File operations might fail.")
             ErrorPopup(message="Storage permission denied.\nThe app might not be able to read/search files correctly.").open()
        else:
             print("Storage permissions granted.")
             # Optional: If folder wasn't set, maybe prompt user now?
             if not self.resume_folder:
                 Clock.schedule_once(lambda dt: ErrorPopup(message="Permissions granted. Please select a folder using 'Change Folder'.").open(), 0.5)


    # ---^^^--- END ADDED on_start ---^^^---


    def build(self):
        self.title = "Search Tool"

        # ---vvv--- MODIFIED FOLDER INITIALIZATION ---vvv---
        self.resume_folder = None # Start with no folder selected
        self.initial_folder_text = "No folder selected"
        # ---^^^--- END MODIFY ---^^^---

        self.exact_match = False
        Window.clearcolor = get_color_from_hex(THEME['background'])

        main_layout = BoxLayout(
            orientation='vertical',
            padding=15,
            spacing=15
        )

        # Background setup (remains the same)
        with main_layout.canvas.before:
            Color(*get_color_from_hex(THEME['background']))
            self.main_bg_rect = Rectangle(pos=main_layout.pos, size=main_layout.size)
        main_layout.bind(pos=self.update_bg_rect, size=self.update_bg_rect)

        # Title Label (remains the same)
        title_label = ThemedLabel(
            text="Search Tool", font_size='20sp', size_hint_y=None, height=50,
            bold=True, color=get_color_from_hex(THEME['primary'])
        )
        main_layout.add_widget(title_label)

        # Search input layout (remains the same)
        search_layout = BoxLayout(orientation='horizontal', size_hint_y=None, height=50, spacing=10)
        self.search_input = ThemedTextInput(hint_text="Enter search keywords", multiline=False)
        search_layout.add_widget(self.search_input)
        search_button = ThemedButton(text="Search", size_hint_x=0.3)
        search_button.bind(on_release=self.search_resumes)
        search_layout.add_widget(search_button)
        main_layout.add_widget(search_layout)

        # Match mode layout (remains the same)
        match_layout = BoxLayout(orientation='horizontal', size_hint_y=None, height=40, spacing=10)
        match_label = ThemedLabel(text="Match mode:", size_hint_x=0.3, halign='right')
        match_label.bind(size=match_label.setter('text_size'))
        match_layout.add_widget(match_label)
        self.exact_button = ThemedButton(text="Partial Match", size_hint_x=0.7)
        self.exact_button.bind(on_release=self.toggle_match_mode)
        match_layout.add_widget(self.exact_button)
        main_layout.add_widget(match_layout)

        # Boolean operators layout (remains the same)
        operators_layout = GridLayout(cols=3, size_hint_y=None, height=50, spacing=10)
        btn_and = ThemedButton(text="AND")
        btn_or = ThemedButton(text="OR")
        btn_not = ThemedButton(text="NOT")
        btn_and.bind(on_release=lambda x: self.append_operator("AND"))
        btn_or.bind(on_release=lambda x: self.append_operator("OR"))
        btn_not.bind(on_release=lambda x: self.append_operator("NOT"))
        operators_layout.add_widget(btn_and)
        operators_layout.add_widget(btn_or)
        operators_layout.add_widget(btn_not)
        main_layout.add_widget(operators_layout)

        # Folder selection layout
        folder_layout = BoxLayout(orientation='horizontal', size_hint_y=None, height=50, spacing=10)

        # ---vvv--- MODIFIED FOLDER LABEL ---vvv---
        self.folder_label = ThemedLabel(
            text=f"Folder: {self.initial_folder_text}", # Use placeholder
            shorten=True, shorten_from='left',
            halign='left',
            valign='middle'
        )
        # ---^^^--- END MODIFY ---^^^---
        self.folder_label.bind(size=self.folder_label.setter('text_size'))
        folder_layout.add_widget(self.folder_label)

        folder_button = ThemedButton(text="Change Folder", size_hint_x=0.3)
        folder_button.bind(on_release=self.show_folder_chooser) # Binding stays
        folder_layout.add_widget(folder_button)
        main_layout.add_widget(folder_layout)

        # Results area setup (remains the same)
        results_card = BoxLayout(orientation='vertical', padding=10, spacing=5)
        with results_card.canvas.before:
            Color(*get_color_from_hex('#ffffff'))
            self.results_rect = Rectangle(pos=results_card.pos, size=results_card.size) # Store rect for update
        results_card.bind(pos=self.update_card_rect, size=self.update_card_rect) # Use method directly

        results_label = ThemedLabel(
            text="Search Results:", size_hint_y=None, height=30, halign='left',
            font_size='16sp', bold=True
        )
        results_label.bind(size=results_label.setter('text_size'))
        results_card.add_widget(results_label)

        self.results_container = BoxLayout(orientation='vertical')
        scroll_view = ScrollView()
        self.results_layout = BoxLayout(orientation='vertical', size_hint_y=None, spacing=3) # Reduced spacing
        self.results_layout.bind(minimum_height=self.results_layout.setter('height'))
        scroll_view.add_widget(self.results_layout)
        self.results_container.add_widget(scroll_view)
        results_card.add_widget(self.results_container)
        main_layout.add_widget(results_card)

        # ---vvv--- ADDED INITIAL CHECK/WARNING ---vvv---
        Clock.schedule_once(self.check_initial_folder, 1) # Delay slightly
        # ---^^^--- END ADD ---^^^---

        return main_layout

    # ---vvv--- ADDED THIS METHOD ---vvv---
    def check_initial_folder(self, dt):
         if not self.resume_folder and platform == 'android': # Only show popup on android maybe?
              ErrorPopup(message="Please select a folder to search using 'Change Folder'.").open()
         elif not self.resume_folder:
              print("App started, no folder selected initially.")
    # ---^^^--- END ADDED ---^^^---

    def update_bg_rect(self, instance, value):
        self.main_bg_rect.pos = instance.pos
        self.main_bg_rect.size = instance.size

    # Modified to use stored rectangle instance directly
    def update_card_rect(self, instance, value):
        if hasattr(self, 'results_rect'):
             self.results_rect.pos = instance.pos
             self.results_rect.size = instance.size

    def append_operator(self, op):
        current_text = self.search_input.text
        # Add space intelligently
        if current_text and not current_text.endswith(' '):
             self.search_input.text = current_text + f" {op} "
        else:
             self.search_input.text = current_text + f"{op} "
        self.search_input.focus = True # Keep focus

    # ---vvv--- REPLACED show_folder_chooser METHOD ---vvv---
    def show_folder_chooser(self, instance):
        """
        Uses plyer.filechooser to open the native directory chooser.
        """
        print("Attempting to show native folder chooser...")
        try:
            # Use choose_dir for selecting directories
            filechooser.choose_dir(on_selection=self.handle_folder_selection)
        except Exception as e:
             print(f"Error opening folder chooser: {e}")
             ErrorPopup(message=f"Could not open folder chooser:\n{e}").open()
    # ---^^^--- END REPLACE ---^^^---

    # ---vvv--- REPLACED set_resume_folder WITH handle_folder_selection ---vvv---
    def handle_folder_selection(self, selection):
        """
        Callback function for plyer's choose_dir. Updates the selected folder.
        """
        print(f"Folder selection callback: {selection}") # selection is a list
        # Plyer often returns a list, even for single selection
        if selection and isinstance(selection, list) and len(selection) > 0:
            selected_path = selection[0]
            # Basic check if it looks like a valid path (might need refinement on Android)
            # os.path.isdir might fail on Android for some SAF paths initially. Trust plyer for now.
            # A better check might involve trying os.listdir shortly after.
            # Let's proceed optimistically first.
            self.resume_folder = selected_path
            # Update label - show basename for readability
            self.folder_label.text = f"Folder: {os.path.basename(self.resume_folder)}"
            # Clear results when folder changes
            self.results_layout.clear_widgets()
            self.results_layout.add_widget(ThemedLabel(text="Folder selected. Ready to search.", size_hint_y=None, height=30))
            print(f"Selected folder set to: {self.resume_folder}")
            # Try a quick listdir to potentially catch immediate permission issues
            try:
                test_list = os.listdir(self.resume_folder)
                print(f"Successfully listed {len(test_list)} items in selected folder.")
            except Exception as list_err:
                 print(f"Warning: Could not immediately list contents of {self.resume_folder}: {list_err}")
                 ErrorPopup(message=f"Warning: Could not list folder contents immediately.\nEnsure permissions are granted.\nError: {list_err}").open()

        else:
            print("Folder selection cancelled or invalid.")
            # Keep self.resume_folder as it was (or None)
            current_display = self.initial_folder_text if not self.resume_folder else os.path.basename(self.resume_folder)
            self.folder_label.text = f"Folder: {current_display}"
    # ---^^^--- END REPLACE ---^^^---

    def toggle_match_mode(self, instance):
        self.exact_match = not self.exact_match
        instance.text = "Exact Match" if self.exact_match else "Partial Match"

    # ---vvv--- MODIFIED search_resumes METHOD ---vvv---
    def search_resumes(self, instance=None): # Allow calling via Enter key
        # --- ADD FOLDER CHECK ---
        if not self.resume_folder: # Check if None or empty string
             ErrorPopup(message="Please select a folder first using 'Change Folder'.").open()
             return
        # Add a check if the directory still exists (could be deleted externally)
        if not os.path.isdir(self.resume_folder):
              ErrorPopup(message=f"Selected folder no longer exists or is invalid:\n{os.path.basename(self.resume_folder)}\nPlease select again.").open()
              # Reset folder state maybe?
              self.resume_folder = None
              self.folder_label.text = f"Folder: {self.initial_folder_text}"
              return
        # --- END FOLDER CHECK ---

        query = self.search_input.text.strip() # Strip whitespace
        if not query:
            ErrorPopup(message="Please enter a search query.").open()
            return

        self.results_layout.clear_widgets()

        status_label = ThemedLabel(
            text=f"Searching in '{os.path.basename(self.resume_folder)}' ({'exact' if self.exact_match else 'partial'})...",
            size_hint_y=None,
            height=30,
            italic=True
        )
        self.results_layout.add_widget(status_label)

        Clock.schedule_once(partial(self.perform_search, query), 0.1)
    # ---^^^--- END MODIFY ---^^^---

    # ---vvv--- MODIFIED perform_search METHOD (Error Handling) ---vvv---
    def perform_search(self, query, dt):
        matching_files = []
        search_error = None # Variable to store potential error message

        # --- ADD CHECK FOR FOLDER EXISTENCE AGAIN ---
        if not self.resume_folder or not os.path.isdir(self.resume_folder):
            search_error = f"Selected folder no longer exists or is invalid:\n{os.path.basename(self.resume_folder) if self.resume_folder else 'None'}"
        else:
            current_folder = self.resume_folder # Cache it in case it changes mid-search? Unlikely but safer.
            try:
                # --- ADD PERMISSION ERROR HANDLING ---
                try:
                    list_of_files = os.listdir(current_folder)
                except PermissionError:
                    search_error = f"Permission denied to read folder:\n{os.path.basename(current_folder)}\nPlease grant storage access and select folder again."
                    list_of_files = []
                except FileNotFoundError: # Handle case where folder disappears between check and listdir
                    search_error = f"Folder not found during search:\n{os.path.basename(current_folder)}"
                    list_of_files = []
                except Exception as list_e: # Catch other listing errors
                     search_error = f"Error listing folder contents:\n{list_e}"
                     list_of_files = []


                if not search_error: # Proceed only if listing succeeded
                     print(f"Found {len(list_of_files)} items in {os.path.basename(current_folder)}") # Debug
                     for file in list_of_files:
                         filepath = os.path.join(current_folder, file)
                         try:
                             # Check if it's a file *and* if we can read it
                             if os.path.isfile(filepath) and os.access(filepath, os.R_OK):
                                 print(f"Processing: {file}") # Debug
                                 text = extract_text(filepath)
                                 if text is not None: # Check if extraction returned something (not None on error)
                                     if boolean_search(text, query, self.exact_match):
                                         matching_files.append(file)
                             # else: # Optional: log skipped items
                             #    if not os.path.isfile(filepath): print(f"Skipping non-file: {file}")
                             #    elif not os.access(filepath, os.R_OK): print(f"Skipping non-readable: {file}")

                         except Exception as process_e:
                              print(f"Error processing file {file}: {process_e}")
                              # Decide whether to stop search or just skip the file
                              # search_error = f"Error processing file {file}: {process_e}" # Example: store last error
                              # break # Example: stop on first file error

            except Exception as e:
                # Catch other potential errors during search setup
                import traceback
                traceback.print_exc()
                search_error = f"Unexpected error during search setup:\n{e}"
        # --- END MODIFICATIONS ---


        # --- UI UPDATE (Clear status, show error or results) ---
        self.results_layout.clear_widgets() # Clear status label

        if search_error:
            ErrorPopup(message=search_error).open()
            no_results = ThemedLabel(
                text="Search failed. See error message.", # Generic message
                size_hint_y=None,
                height=40,
                color=get_color_from_hex(THEME['accent'])
            )
            self.results_layout.add_widget(no_results)
            return # Stop processing results

        if matching_files:
            result_count = ThemedLabel(
                text=f"Found {len(matching_files)} matching file(s):",
                size_hint_y=None,
                height=40,
                halign='left',
                color=get_color_from_hex(THEME['success'])
            )
            result_count.bind(size=result_count.setter('text_size'))
            self.results_layout.add_widget(result_count)

            matching_files.sort() # Sort results alphabetically

            for filename in matching_files:
                result_item = ResultItem(filename=filename) # Pass filename correctly
                self.results_layout.add_widget(result_item)
        else:
            no_results = ThemedLabel(
                text="No matching files found.",
                size_hint_y=None,
                height=40,
                color=get_color_from_hex(THEME['text']) # Use normal text color for 'not found'
            )
            self.results_layout.add_widget(no_results)
        # --- END UI UPDATE ---

    # ---^^^--- END MODIFY perform_search ---^^^---


if __name__ == "__main__":
    # Ensure necessary dirs exist? Not usually needed for App().user_data_dir
    ResumeSearchApp().run()


# --- END OF FILE SearchTool.py ---
