import os
import re
import fitz  # PyMuPDF for PDFs
import docx
import pptx
import pandas as pd
import tkinter as tk
from tkinter import filedialog, messagebox, scrolledtext, Menu
import subprocess
import platform

RESUME_FOLDER = "Resume_Download"

def extract_text_from_pdf(filepath):
    text = ""
    try:
        doc = fitz.open(filepath)
        for page in doc:
            text += page.get_text("text") + "\n"
    except Exception as e:
        print(f"Error reading {filepath}: {e}")
    return text

def extract_text_from_docx(filepath):
    try:
        doc = docx.Document(filepath)
        return "\n".join([para.text for para in doc.paragraphs])
    except Exception as e:
        print(f"Error reading {filepath}: {e}")
        return ""

def extract_text_from_pptx(filepath):
    try:
        presentation = pptx.Presentation(filepath)
        text = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)
    except Exception as e:
        print(f"Error reading {filepath}: {e}")
        return ""

def extract_text_from_excel(filepath):
    try:
        df = pd.read_excel(filepath, sheet_name=None)
        text = "\n".join([df[sheet].to_string() for sheet in df])
        return text
    except Exception as e:
        print(f"Error reading {filepath}: {e}")
        return ""

def extract_text(filepath):
    ext = filepath.split(".")[-1].lower()
    if ext == "pdf":
        return extract_text_from_pdf(filepath)
    elif ext == "docx":
        return extract_text_from_docx(filepath)
    elif ext == "pptx":
        return extract_text_from_pptx(filepath)
    elif ext in ["xls", "xlsx"]:
        return extract_text_from_excel(filepath)
    else:
        print(f"Unsupported format: {filepath}")
        return ""

def boolean_search(text, query):
    query = query.replace("AND", "&").replace("OR", "|").replace("NOT", "~")
    words = re.findall(r'\w+', query)
    
    for word in words:
        if word.lower() not in text.lower():
            query = query.replace(word, "False")
        else:
            query = query.replace(word, "True")
    
    try:
        return eval(query)
    except:
        return False

def on_select(event):
    try:
        result_text.tag_remove("highlight", "1.0", tk.END)  # Remove old selection
        index = result_text.index(tk.CURRENT)
        line_start = f"{index.split('.')[0]}.0"
        line_end = f"{index.split('.')[0]}.end"

        result_text.tag_add("highlight", line_start, line_end)  # Add highlight tag
        result_text.tag_config("highlight", background="lightblue")  # Blue highlight
    except Exception as e:
        print(f"Error selecting file: {e}")

def open_file(event):
    try:
        # Find the selected line based on highlight
        index = result_text.index(tk.CURRENT)
        line_start = f"{index.split('.')[0]}.0"
        line_end = f"{index.split('.')[0]}.end"
        selected_text = result_text.get(line_start, line_end).strip()

        if selected_text:
            filepath = os.path.join(RESUME_FOLDER, selected_text)
            if os.path.exists(filepath):
                if platform.system() == "Windows":
                    os.startfile(filepath)
                elif platform.system() == "Darwin":  # macOS
                    subprocess.run(["open", filepath])
                else:  # Linux
                    subprocess.run(["xdg-open", filepath])
    except Exception as e:
        print(f"Error opening file: {e}")

def search_resumes():
    query = search_entry.get()
    if not query:
        messagebox.showerror("Error", "Please enter a search query.")
        return
    
    matching_files = []
    for file in os.listdir(RESUME_FOLDER):
        filepath = os.path.join(RESUME_FOLDER, file)
        if os.path.isfile(filepath):
            text = extract_text(filepath)
            if boolean_search(text, query):
                matching_files.append(file)
    
    result_text.delete("1.0", tk.END)
    if matching_files:
        result_text.insert(tk.END, "\n".join(matching_files))
    else:
        result_text.insert(tk.END, "No matching resumes found.")

def append_operator(op):
    text = search_entry.get()
    search_entry.delete(0, tk.END)
    search_entry.insert(0, text + f" {op} ")

def browse_folder():
    global RESUME_FOLDER
    folder_selected = filedialog.askdirectory()
    if folder_selected:
        RESUME_FOLDER = folder_selected
        folder_label.config(text=f"Folder: {RESUME_FOLDER}")

# GUI Setup
root = tk.Tk()
root.title("Search String")
root.geometry("500x400")

try:
    root.iconbitmap(r".\icon1.ico")
except tk.TclError:
    print(f"Icon file '{icon_file_pathy}' not found or not a valid icon file.")

tk.Label(root, text="Enter search keywords:").pack(pady=5)
search_entry = tk.Entry(root, width=50)
search_entry.pack(pady=5)

button_frame = tk.Frame(root)
button_frame.pack(pady=5)

tk.Button(button_frame, text="AND", command=lambda: append_operator("AND")).pack(side=tk.LEFT, padx=5)
tk.Button(button_frame, text="OR", command=lambda: append_operator("OR")).pack(side=tk.LEFT, padx=5)
tk.Button(button_frame, text="NOT", command=lambda: append_operator("NOT")).pack(side=tk.LEFT, padx=5)

tk.Button(root, text="Search", command=search_resumes).pack(pady=5)
folder_label = tk.Label(root, text=f"Folder: {RESUME_FOLDER}")
folder_label.pack(pady=5)
tk.Button(root, text="Change Folder", command=browse_folder).pack(pady=5)

result_text = scrolledtext.ScrolledText(root, height=10, width=60)
result_text.pack(pady=5)
result_text.bind("<Button-1>", on_select)  # Left-click to highlight
result_text.bind("<Button-3>", open_file)  # Right-click to open file

root.mainloop()
