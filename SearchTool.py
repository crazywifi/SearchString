import os
import re
import fitz  # PyMuPDF for PDFs
import docx
import pptx
import pandas as pd
import subprocess
import platform
from kivy.app import App
from kivy.uix.boxlayout import BoxLayout
from kivy.uix.gridlayout import GridLayout
from kivy.uix.button import Button
from kivy.uix.label import Label
from kivy.uix.textinput import TextInput
from kivy.uix.scrollview import ScrollView
from kivy.uix.filechooser import FileChooserListView
from kivy.uix.popup import Popup
from kivy.core.window import Window
from kivy.uix.behaviors import ButtonBehavior
from kivy.uix.modalview import ModalView
from kivy.graphics import Color, Rectangle
from kivy.utils import get_color_from_hex
from kivy.clock import Clock
from functools import partial

# Set theme colors
THEME = {
    'primary': '#738c63',  # Blue
    'secondary': '#A0B38C',  # Light blue
    'text': '#333333',  # Dark gray
    'background': '#A0B38C',  # Off-white
    'accent': '#ff6b6b',  # Accent red
    'success': '#5cb85c',  # Green
    'hover': '#d9e6f7'  # Hover light blue
}
RESULT_ITEM_COLORS = {
    'normal': '#b0c49a',  # Lighter Green (Slightly lighter than background)
    'hover': '#A0B38C',   # Main Background Green (Theme secondary)
    'press': '#889877'    # Darker Green
}

#RESUME_FOLDER = "Resume_Download"

def extract_text_from_pdf(filepath):
    text = ""
    try:
        doc = fitz.open(filepath)
        for page in doc:
            text += page.get_text("text") + "\n"
    except Exception as e:
        print(f"Error reading {filepath}: {e}")
    return text

def extract_text_from_docx(filepath):
    try:
        doc = docx.Document(filepath)
        return "\n".join([para.text for para in doc.paragraphs])
    except Exception as e:
        print(f"Error reading {filepath}: {e}")
        return ""

def extract_text_from_pptx(filepath):
    try:
        presentation = pptx.Presentation(filepath)
        text = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)
    except Exception as e:
        print(f"Error reading {filepath}: {e}")
        return ""

def extract_text_from_excel(filepath):
    try:
        df = pd.read_excel(filepath, sheet_name=None)
        text = "\n".join([df[sheet].to_string() for sheet in df])
        return text
    except Exception as e:
        print(f"Error reading {filepath}: {e}")
        return ""

def extract_text_from_txt(filepath):
    try:
        with open(filepath, 'r', encoding='utf-8', errors='ignore') as file:
            return file.read()
    except Exception as e:
        print(f"Error reading {filepath}: {e}")
        return ""

def extract_text_from_csv(filepath):
    try:
        df = pd.read_csv(filepath)
        return df.to_string()
    except Exception as e:
        print(f"Error reading {filepath}: {e}")
        return ""

def extract_text(filepath):
    ext = filepath.split(".")[-1].lower()
    if ext == "pdf":
        return extract_text_from_pdf(filepath)
    elif ext == "docx":
        return extract_text_from_docx(filepath)
    elif ext == "pptx":
        return extract_text_from_pptx(filepath)
    elif ext in ["xls", "xlsx"]:
        return extract_text_from_excel(filepath)
    elif ext == "txt":
        return extract_text_from_txt(filepath)
    elif ext == "csv":
        return extract_text_from_csv(filepath)
    elif ext == "rtf":
        try:
            with open(filepath, 'r', encoding='utf-8', errors='ignore') as file:
                return file.read()
        except Exception as e:
            print(f"Error reading {filepath}: {e}")
            return ""
    elif ext in ["json", "xml", "html", "htm", "md", "log"]:
        return extract_text_from_txt(filepath)
    else:
        print(f"Unsupported format: {filepath}")
        return ""

def boolean_search(text, query, exact_match=False):
    query = query.replace("AND", "&").replace("OR", "|").replace("NOT", "~")
    words = re.findall(r'\w+', query)
    
    for word in words:
        if exact_match:
            # For exact match, look for word boundaries
            pattern = r'\b' + re.escape(word) + r'\b'
            if not re.search(pattern, text, re.IGNORECASE):
                query = query.replace(word, "False")
            else:
                query = query.replace(word, "True")
        else:
            # For partial match
            if word.lower() not in text.lower():
                query = query.replace(word, "False")
            else:
                query = query.replace(word, "True")
    
    try:
        return eval(query)
    except:
        return False

class ThemedButton(Button):
    def __init__(self, **kwargs):
        super(ThemedButton, self).__init__(**kwargs)
        self.background_normal = ''
        self.background_color = get_color_from_hex(THEME['primary'])
        self.color = get_color_from_hex('#ffffff')  # White text
        self.border = (0, 0, 0, 0)
        
    def on_press(self):
        self.background_color = get_color_from_hex('#395582')  # Darker when pressed
        
    def on_release(self):
        self.background_color = get_color_from_hex(THEME['primary'])

class ThemedLabel(Label):
    def __init__(self, **kwargs):
        super(ThemedLabel, self).__init__(**kwargs)
        self.color = get_color_from_hex(THEME['text'])

class ThemedTextInput(TextInput):
    def __init__(self, **kwargs):
        super(ThemedTextInput, self).__init__(**kwargs)
        self.background_color = get_color_from_hex('#ffffff')
        self.foreground_color = get_color_from_hex(THEME['text'])
        self.cursor_color = get_color_from_hex(THEME['primary'])
        self.selection_color = get_color_from_hex(THEME['secondary'])
        self.padding = [10, 10, 10, 10]
        self.cursor_width = 2

class ResultItem(ButtonBehavior, BoxLayout):
    def __init__(self, filename, **kwargs):
        super(ResultItem, self).__init__(**kwargs)
        self.orientation = 'horizontal'
        self.size_hint_y = None
        self.height = 50
        self.filename = filename
        # self.click_time = 0 # Original line, can be kept or removed if not used
        self.last_click = 0
        self.is_hovering = False # ADDED: Track hover state

        # ADDED: Define custom colors for different states using RESULT_ITEM_COLORS
        self.normal_color_hex = RESULT_ITEM_COLORS['normal']
        self.hover_color_hex = RESULT_ITEM_COLORS['hover']
        self.press_color_hex = RESULT_ITEM_COLORS['press']

        # ADDED: Initial background color setup using canvas
        with self.canvas.before:
            self.bg_color = Color(*get_color_from_hex(self.normal_color_hex))
            self.bg_rect = Rectangle(pos=self.pos, size=self.size)

        # MODIFIED: Bind position and size to update background rectangle
        self.bind(pos=self.update_rect, size=self.update_rect)

        # Label setup (adjust as per your original if needed, ensure text shortening)
        self.label = ThemedLabel(
            text=filename,
            size_hint_x=0.9,
            halign='left',
            valign='middle',
            shorten=True, # Ensure text shortening if needed
            shorten_from='right'
        )
        self.label.bind(size=self.label.setter('text_size'))
        self.add_widget(self.label)

        # MODIFIED: Bind ButtonBehavior events (ensure these are present)
        self.bind(on_press=self.on_item_press)
        self.bind(on_release=self.on_item_release)

        # ADDED: Bind mouse hover events using Window
        Window.bind(mouse_pos=self.on_mouse_pos) # Bind window mouse pos changes

    # MODIFIED: Ensure update_rect updates the background and label text_size
    def update_rect(self, *args):
        self.bg_rect.pos = self.pos
        self.bg_rect.size = self.size
        # Optional but good: update label text size constraint
        if hasattr(self, 'label'):
             self.label.text_size = (self.width * 0.9, None)

    # --- ADDED: New methods for hover handling ---
    def on_mouse_pos(self, window, pos):
        # Check if the mouse position is within the bounds of this widget
        # Need to convert window coords to local widget coords
        widget_pos = self.to_widget(*pos)
        is_inside = self.collide_point(*widget_pos)

        if self.is_hovering != is_inside:
            # Hover state changed
            self.is_hovering = is_inside
            if is_inside:
                # Check state to avoid changing color if button is pressed down
                if self.state == 'normal':
                    self.on_enter()
            else:
                # Check state to avoid changing color if button is pressed down
                if self.state == 'normal':
                    self.on_leave()

    def on_enter(self):
        """Called when mouse enters the widget area."""
        Window.set_system_cursor('hand')
        self.bg_color.rgba = get_color_from_hex(self.hover_color_hex)

    def on_leave(self):
        """Called when mouse leaves the widget area."""
        Window.set_system_cursor('arrow')
        self.bg_color.rgba = get_color_from_hex(self.normal_color_hex)
    # --- End of Added Hover Methods ---

    # MODIFIED: Update on_item_press to set pressed color
    def on_item_press(self, instance):
        """Called when the item is pressed."""
        self.bg_color.rgba = get_color_from_hex(self.press_color_hex) # Set press color
        current_time = Clock.get_time()

        # Check for double click (keep your original logic)
        if current_time - self.last_click < 0.3:
            self.open_file()

        self.last_click = current_time

    # MODIFIED: Update on_item_release for correct hover/normal state after release
    def on_item_release(self, instance):
        """Called when the item is released."""
        # Check if mouse is still over the widget upon release
        # Need to re-check position in case mouse moved *while* pressed
        current_mouse_pos = Window.mouse_pos
        widget_pos = self.to_widget(*current_mouse_pos)
        is_still_inside = self.collide_point(*widget_pos)

        if is_still_inside:
            self.is_hovering = True # Ensure hover state is correct
            self.bg_color.rgba = get_color_from_hex(self.hover_color_hex)
            Window.set_system_cursor('hand') # Ensure cursor is hand if still hovering
        else:
            self.is_hovering = False # Ensure hover state is correct
            # If mouse moved off *while* pressed, revert to normal
            self.bg_color.rgba = get_color_from_hex(self.normal_color_hex)
            Window.set_system_cursor('arrow')

    # Keep your original open_file method
    def open_file(self):
        try:
            # Use the currently stored resume_folder from the App instance
            app = App.get_running_app()
            if not app or not hasattr(app, 'resume_folder'):
                 print("Error: App instance or resume_folder not found.")
                 # Maybe show an error popup here?
                 # ErrorPopup(message="Internal error: Cannot find app folder reference.").open()
                 return

            filepath = os.path.join(app.resume_folder, self.filename)
            if os.path.exists(filepath):
                if platform.system() == "Windows":
                    os.startfile(filepath)
                elif platform.system() == "Darwin":  # macOS
                    subprocess.run(["open", filepath], check=True)
                else:  # Linux
                    subprocess.run(["xdg-open", filepath], check=True)
            else:
                print(f"Error: File not found at {filepath}")
                ErrorPopup(message=f"File not found:\n{self.filename}").open()
        except Exception as e:
            print(f"Error opening file '{self.filename}': {e}")
            ErrorPopup(message=f"Could not open file:\n{e}").open()


class CustomFileChooserListView(FileChooserListView):
    def __init__(self, **kwargs):
        super(CustomFileChooserListView, self).__init__(**kwargs)
        # Set text color for entries
        self.layout.ids.entries.color = get_color_from_hex(THEME['text'])


        
class FolderChooserPopup(Popup):
    def __init__(self, callback, **kwargs):
        super(FolderChooserPopup, self).__init__(**kwargs)
        self.title = "Select Resume Folder"
        self.size_hint = (0.9, 0.9)
        self.callback = callback
        self.background = ''
        self.background_color = get_color_from_hex('#ffffff')
        self.title_color = get_color_from_hex(THEME['primary'])
        self.separator_color = get_color_from_hex(THEME['primary'])
        
        layout = BoxLayout(orientation='vertical', spacing=10, padding=10)
        
        # Create a BoxLayout with a custom background for the file chooser
        file_chooser_container = BoxLayout(orientation='vertical')
        with file_chooser_container.canvas.before:
            Color(*get_color_from_hex('#f5f5f5'))  # Light gray background
            Rectangle(pos=file_chooser_container.pos, size=file_chooser_container.size)
        file_chooser_container.bind(pos=self.update_container_rect, size=self.update_container_rect)
        
        self.file_chooser = FileChooserListView(
            path=os.getcwd(), 
            dirselect=True,
            file_encodings=['utf-8', 'latin1', 'cp1252']
        )
        
        # Customize the file chooser appearance
        self.file_chooser.background_color = get_color_from_hex('#738c63')
        self.file_chooser.background_normal = ''
        self.file_chooser.foreground_color = get_color_from_hex(THEME['text'])  # Dark text color
        
        # Add the file chooser to the container
        file_chooser_container.add_widget(self.file_chooser)
        
        button_layout = BoxLayout(size_hint_y=None, height=50, spacing=10)
        
        btn_cancel = ThemedButton(text="Cancel")
        btn_select = ThemedButton(text="Select")
        
        btn_cancel.bind(on_release=self.dismiss)
        btn_select.bind(on_release=self._select_folder)
        
        button_layout.add_widget(btn_cancel)
        button_layout.add_widget(btn_select)
        
        layout.add_widget(file_chooser_container)
        layout.add_widget(button_layout)
        
        self.content = layout
    
    def update_container_rect(self, instance, value):
        # Update the background rectangle when the container size changes
        instance.canvas.before.clear()
        with instance.canvas.before:
            Color(*get_color_from_hex('#738c63'))  # Light gray background
            Rectangle(pos=instance.pos, size=instance.size)
        
    def _select_folder(self, instance):
        selected = self.file_chooser.path
        self.callback(selected)
        self.dismiss()

class ErrorPopup(Popup):
    def __init__(self, message, **kwargs):
        super(ErrorPopup, self).__init__(**kwargs)
        self.title = "Error"
        self.size_hint = (0.8, 0.3)
        self.auto_dismiss = True
        self.background = ''
        self.background_color = get_color_from_hex('#ffffff')
        self.title_color = get_color_from_hex(THEME['accent'])
        self.separator_color = get_color_from_hex(THEME['accent'])
        
        content = BoxLayout(orientation='vertical', spacing=10, padding=10)
        content.add_widget(ThemedLabel(text=message))
        
        btn = ThemedButton(text="OK", size_hint_y=None, height=50)
        btn.bind(on_release=self.dismiss)
        content.add_widget(btn)
        
        self.content = content

class ResumeSearchApp(App):
    def build(self):
        self.title = "Search Tool"
        #self.resume_folder = RESUME_FOLDER
        #self.exact_match = False

        try:
            script_dir = os.path.dirname(os.path.abspath(__file__))
        except NameError:
             # Fallback if __file__ is not defined (e.g., interactive mode)
             script_dir = os.getcwd()
        self.resume_folder = script_dir # Set this as the default folder
        # ---^^^--- END OF REPLACEMENT/ADDITION ---^^^---

        self.exact_match = False
        # Set window background color
        Window.clearcolor = get_color_from_hex(THEME['background'])
        
        # Main layout
        main_layout = BoxLayout(
            orientation='vertical', 
            padding=15, 
            spacing=15
        )
        
        # Set background color for main layout using canvas
        with main_layout.canvas.before:
            Color(*get_color_from_hex(THEME['background']))
            self.main_bg_rect = Rectangle(pos=main_layout.pos, size=main_layout.size)
        main_layout.bind(pos=self.update_bg_rect, size=self.update_bg_rect)
        
        # Add title
        title_label = ThemedLabel(
            text="Search Tool",
            font_size='20sp',
            size_hint_y=None,
            height=50,
            bold=True,
            color=get_color_from_hex(THEME['primary'])
        )
        main_layout.add_widget(title_label)
        
        # Search input
        search_layout = BoxLayout(orientation='horizontal', size_hint_y=None, height=50, spacing=10)
        
        self.search_input = ThemedTextInput(
            hint_text="Enter search keywords",
            multiline=False
        )
        search_layout.add_widget(self.search_input)
        
        search_button = ThemedButton(text="Search", size_hint_x=0.3)
        search_button.bind(on_release=self.search_resumes)
        search_layout.add_widget(search_button)
        
        main_layout.add_widget(search_layout)
        
        # Exact match checkbox
        match_layout = BoxLayout(orientation='horizontal', size_hint_y=None, height=40, spacing=10)
        
        match_label = ThemedLabel(text="Match mode:", size_hint_x=0.3, halign='right')
        match_label.bind(size=match_label.setter('text_size'))
        match_layout.add_widget(match_label)
        
        self.exact_button = ThemedButton(text="Partial Match", size_hint_x=0.7)
        self.exact_button.bind(on_release=self.toggle_match_mode)
        match_layout.add_widget(self.exact_button)
        
        main_layout.add_widget(match_layout)
        
        # Boolean operators
        operators_layout = GridLayout(cols=3, size_hint_y=None, height=50, spacing=10)
        
        btn_and = ThemedButton(text="AND")
        btn_or = ThemedButton(text="OR")
        btn_not = ThemedButton(text="NOT")
        
        btn_and.bind(on_release=lambda x: self.append_operator("AND"))
        btn_or.bind(on_release=lambda x: self.append_operator("OR"))
        btn_not.bind(on_release=lambda x: self.append_operator("NOT"))
        
        operators_layout.add_widget(btn_and)
        operators_layout.add_widget(btn_or)
        operators_layout.add_widget(btn_not)
        
        main_layout.add_widget(operators_layout)
        
        # Folder selection
        folder_layout = BoxLayout(orientation='horizontal', size_hint_y=None, height=50, spacing=10)
        
        self.folder_label = ThemedLabel(
            text=f"Folder: {self.resume_folder}",
            valign='middle'
        )
        self.folder_label.bind(size=self.folder_label.setter('text_size'))
        folder_layout.add_widget(self.folder_label)
        
        folder_button = ThemedButton(text="Change Folder", size_hint_x=0.3)
        folder_button.bind(on_release=self.show_folder_chooser)
        folder_layout.add_widget(folder_button)
        
        main_layout.add_widget(folder_layout)
        
        # Results area
        results_card = BoxLayout(
            orientation='vertical',
            padding=10,
            spacing=5
        )
        
        with results_card.canvas.before:
            Color(*get_color_from_hex('#ffffff'))
            Rectangle(pos=results_card.pos, size=results_card.size)
        
        results_card.bind(pos=self.update_card_rect, size=self.update_card_rect)
        
        results_label = ThemedLabel(
            text="Search Results:",
            size_hint_y=None, 
            height=30, 
            halign='left',
            font_size='16sp',
            bold=True
        )
        results_label.bind(size=results_label.setter('text_size'))
        results_card.add_widget(results_label)
        
        self.results_container = BoxLayout(orientation='vertical')
        
        scroll_view = ScrollView()
        self.results_layout = BoxLayout(orientation='vertical', size_hint_y=None, spacing=5)
        self.results_layout.bind(minimum_height=self.results_layout.setter('height'))
        
        scroll_view.add_widget(self.results_layout)
        self.results_container.add_widget(scroll_view)
        
        results_card.add_widget(self.results_container)
        main_layout.add_widget(results_card)
        
        # Store the card for updates
        self.results_card = results_card
        
        return main_layout
    
    def update_bg_rect(self, instance, value):
        self.main_bg_rect.pos = instance.pos
        self.main_bg_rect.size = instance.size
    
    def update_card_rect(self, instance, value):
        instance.canvas.before.clear()
        with instance.canvas.before:
            Color(*get_color_from_hex('#ffffff'))
            Rectangle(pos=instance.pos, size=instance.size)
    
    def append_operator(self, op):
        current_text = self.search_input.text
        self.search_input.text = current_text + f" {op} "
    
    def show_folder_chooser(self, instance):
        popup = FolderChooserPopup(callback=self.set_resume_folder)
        popup.open()
    
    def set_resume_folder(self, folder_path):
        self.resume_folder = folder_path
        self.folder_label.text = f"Folder: {self.resume_folder}"
    
    def toggle_match_mode(self, instance):
        self.exact_match = not self.exact_match
        if self.exact_match:
            self.exact_button.text = "Exact Match"
        else:
            self.exact_button.text = "Partial Match"
    
    def search_resumes(self, instance):
        query = self.search_input.text
        if not query:
            ErrorPopup(message="Please enter a search query.").open()
            return
        
        # Clear previous results
        self.results_layout.clear_widgets()
        
        # Add a status label
        status_label = ThemedLabel(
            text=f"Searching with {'exact' if self.exact_match else 'partial'} matching...", 
            size_hint_y=None, 
            height=30
        )
        self.results_layout.add_widget(status_label)
        
        # Use Clock to prevent UI from freezing during search
        Clock.schedule_once(partial(self.perform_search, query), 0.1)
    
    def perform_search(self, query, dt):
        matching_files = []
        try:
            if os.path.exists(self.resume_folder):
                for file in os.listdir(self.resume_folder):
                    filepath = os.path.join(self.resume_folder, file)
                    if os.path.isfile(filepath):
                        text = extract_text(filepath)
                        if boolean_search(text, query, self.exact_match):
                            matching_files.append(file)
            else:
                raise FileNotFoundError(f"Folder not found: {self.resume_folder}")
        except Exception as e:
            self.results_layout.clear_widgets()
            ErrorPopup(message=f"Error searching files: {str(e)}").open()
            return
        
        # Clear the status label
        self.results_layout.clear_widgets()
        
        if matching_files:
            result_count = ThemedLabel(
                text=f"Found {len(matching_files)} matching file(s)", 
                size_hint_y=None, 
                height=40,
                halign='left',
                color=get_color_from_hex(THEME['success'])
            )
            result_count.bind(size=result_count.setter('text_size'))
            self.results_layout.add_widget(result_count)
            
            for filename in matching_files:
                result_item = ResultItem(filename)
                self.results_layout.add_widget(result_item)
        else:
            no_results = ThemedLabel(
                text="No matching files found.", 
                size_hint_y=None, 
                height=40,
                color=get_color_from_hex(THEME['accent'])
            )
            self.results_layout.add_widget(no_results)

if __name__ == "__main__":
    ResumeSearchApp().run()
